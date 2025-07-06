import streamlit as st
import pytesseract
from PIL import Image
from deep_translator import GoogleTranslator
from gtts import gTTS
import io
import firebase_admin
from firebase_admin import credentials, db
import re
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation

# ---------- ตั้งค่าหน้าเว็บ ----------
st.set_page_config(page_title="📘 Word & Sentence Translator", layout="centered")
st.title("📘 โปรแกรมแปลศัพท์ + ประโยค + เสียงอ่าน")

# ---------- ฟังก์ชัน Normalize ----------
def normalize(word):
    return re.sub(r"[^a-zA-Z0-9\-]", "", word).strip().lower()

# ---------- Firebase Init ----------
def init_firebase():
    if not firebase_admin._apps:
        firebase_config = st.secrets["FIREBASE"]
        cred = credentials.Certificate(firebase_config)
        firebase_admin.initialize_app(cred, {
            'databaseURL': 'https://vocab-tracker-7e059-default-rtdb.asia-southeast1.firebasedatabase.app/'
        })

# ---------- บันทึกคำศัพท์ใหม่ ----------
def save_to_firebase(data):
    init_firebase()
    ref = db.reference('vocabulary')
    existing_data = ref.get()
    existing_words = set()

    if existing_data:
        for item in existing_data.values():
            existing_words.add(normalize(item.get("english", "")))

    added_count = 0
    for word, translation in data:
        cleaned = normalize(word)
        if cleaned and cleaned not in existing_words:
            ref.push({
                "english": word.strip(),
                "thai": translation.strip()
            })
            existing_words.add(cleaned)
            added_count += 1

    if added_count == 0:
        st.info("📌 ไม่มีคำใหม่เพิ่ม เพราะคำศัพท์ทั้งหมดมีอยู่แล้วใน Firebase")
    else:
        st.success(f"✅ บันทึกคำใหม่ {added_count} คำ ลง Firebase เรียบร้อยแล้ว")

# ---------- อัปโหลดไฟล์ ----------
uploaded_file = st.file_uploader("📤 อัปโหลดไฟล์ (.jpg, .png, .pdf, .pptx)", type=["jpg", "jpeg", "png", "pdf", "pptx"])
text = ""

if uploaded_file:
    file_type = uploaded_file.type

    # ---------- IMAGE ----------
    if file_type.startswith("image/"):
        image = Image.open(uploaded_file)
        st.image(image, caption='📷 ภาพต้นฉบับ', use_container_width=True)
        gray_image = image.convert("L")
        text = pytesseract.image_to_string(gray_image, lang="eng")

    # ---------- PDF ----------
    elif file_type == "application/pdf":
        reader = PdfReader(uploaded_file)
        total_pages = len(reader.pages)
        st.write(f"📄 พบทั้งหมด {total_pages} หน้า")

        mode = st.radio("เลือกหน้าที่ต้องการแปล", ["📌 เลือกหน้าเดียว", "📍 ช่วงหน้า", "📚 ทุกหน้า"])
        if mode == "📌 เลือกหน้าเดียว":
            page = st.selectbox("📄 เลือกหน้า", list(range(1, total_pages + 1)), index=0)
            text = reader.pages[page - 1].extract_text()
        elif mode == "📍 ช่วงหน้า":
            start = st.number_input("เริ่มหน้า", 1, total_pages, value=1)
            end = st.number_input("ถึงหน้า", start, total_pages, value=start)
            texts = [reader.pages[i - 1].extract_text() for i in range(start, end + 1)]
            text = "\n\n".join(texts)
        else:
            all_text = [p.extract_text() for p in reader.pages]
            text = "\n\n".join(all_text)

    # ---------- PPTX ----------
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        prs = Presentation(uploaded_file)
        slide_texts = ["\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]) for slide in prs.slides]
        total_slides = len(slide_texts)
        st.write(f"📊 พบทั้งหมด {total_slides} สไลด์")

        mode = st.radio("เลือกสไลด์", ["📌 สไลด์เดียว", "📍 ช่วงสไลด์", "📚 ทุกสไลด์"])
        if mode == "📌 สไลด์เดียว":
            slide = st.selectbox("เลือกสไลด์", list(range(1, total_slides + 1)), index=0)
            text = slide_texts[slide - 1]
        elif mode == "📍 ช่วงสไลด์":
            start = st.number_input("เริ่มสไลด์", 1, total_slides, value=1)
            end = st.number_input("ถึงสไลด์", start, total_slides, value=start)
            text = "\n\n".join(slide_texts[start - 1:end])
        else:
            text = "\n\n".join(slide_texts)

# ---------- ตรวจแก้ข้อความ ----------
if text:
    editable_text = st.text_area("📋 ข้อความที่ตรวจจับได้ (แก้ไขได้)", value=text.strip(), height=200)

    if st.button("🔊 อ่านข้อความภาษาอังกฤษ"):
        tts_en = gTTS(editable_text, lang='en')
        en_audio = io.BytesIO()
        tts_en.write_to_fp(en_audio)
        en_audio.seek(0)
        st.audio(en_audio, format='audio/mp3')

    if st.button("🧠 แปลคำศัพท์"):
        words = []
        for line in editable_text.splitlines():
            for word in line.split():
                clean_word = re.sub(r"[^a-zA-Z0-9\-]", "", word)
                if clean_word:
                    words.append(clean_word)

        vocab = []
        for word in words:
            try:
                th = GoogleTranslator(source="en", target="th").translate(word)
            except:
                th = "-"
            vocab.append({"english": word, "thai": th})

        st.session_state["vocab_df"] = pd.DataFrame(vocab)

# ---------- แก้ไขคำศัพท์ก่อนบันทึก ----------
if "vocab_df" in st.session_state:
    st.info("✏️ แก้ไข/ลบคำศัพท์ก่อนบันทึก")
    edited_df = st.data_editor(
        st.session_state["vocab_df"],
        use_container_width=True,
        num_rows="dynamic",
        hide_index=True,
        column_config={
            "english": st.column_config.TextColumn(label="English"),
            "thai": st.column_config.TextColumn(label="Thai Translation")
        }
    )

    if st.button("💾 บันทึกคำศัพท์ลง Firebase"):
        if not edited_df.empty:
            data_to_save = list(zip(edited_df["english"], edited_df["thai"]))
            save_to_firebase(data_to_save)
            del st.session_state["vocab_df"]
        else:
            st.warning("⚠️ ไม่มีคำศัพท์ที่จะแสดง")

# ---------- แปลเป็นประโยค ----------
if text:
    st.subheader("📝 แปลเป็นประโยคเต็ม")
    try:
        full_translation = GoogleTranslator(source='en', target='th').translate(editable_text)
        st.success(full_translation)

        if st.button("🔊 อ่านคำแปลภาษาไทย"):
            tts_th = gTTS(full_translation, lang='th')
            th_audio = io.BytesIO()
            tts_th.write_to_fp(th_audio)
            th_audio.seek(0)
            st.audio(th_audio, format='audio/mp3')
    except Exception as e:
        st.error("⚠️ แปลประโยคไม่ได้: " + str(e))